# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "python-pptx",
# ]
# ///
"""Helpers for creating a new PPTX presentation from an existing template.

Usage:
    uv run .github/skills/pptx-from-template/pptx_from_template.py template.pptx

All functions operate on python-pptx Presentation objects or file paths.
"""

from __future__ import annotations

import copy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def inspect_template(template_path: str | Path) -> dict[str, Any]:
    """Inspect a PPTX template and return its structure.

    Returns a dict with:
        - dimensions: (width_inches, height_inches)
        - layouts: list of {index, name, placeholders: [{idx, name, type}]}
        - slides: list of {index, layout_name, shapes: [{name, text}]}
    """
    prs = Presentation(str(template_path))

    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": str(ph.placeholder_format.type),
            })
        layouts.append({"index": i, "name": layout.name, "placeholders": placeholders})

    slides = []
    for i, slide in enumerate(prs.slides):
        shapes = []
        for shape in slide.shapes:
            shape_info: dict[str, Any] = {"name": shape.name}
            if shape.has_text_frame:
                shape_info["text"] = shape.text_frame.text[:200]
            shapes.append(shape_info)
        slides.append({
            "index": i + 1,
            "layout_name": slide.slide_layout.name,
            "shapes": shapes,
        })

    return {
        "dimensions": (
            round(prs.slide_width / 914400, 1),
            round(prs.slide_height / 914400, 1),
        ),
        "layouts": layouts,
        "slides": slides,
    }


def print_template_info(template_path: str | Path) -> None:
    """Print a human-readable summary of a template's structure."""
    info = inspect_template(template_path)
    w, h = info["dimensions"]
    print(f"Slide dimensions: {w}\" x {h}\"\n")

    print("Available layouts:")
    for layout in info["layouts"]:
        ph_summary = ", ".join(f"idx={p['idx']} ({p['name']})" for p in layout["placeholders"])
        print(f"  [{layout['index']}] '{layout['name']}': {ph_summary}")

    print(f"\nExisting slides ({len(info['slides'])}):")
    for slide in info["slides"]:
        print(f"  Slide {slide['index']}: layout='{slide['layout_name']}'")
        for shape in slide["shapes"]:
            if "text" in shape:
                preview = shape["text"][:80].replace("\n", " ")
                print(f"    '{shape['name']}': \"{preview}\"")


def _delete_all_slides(prs: Presentation) -> None:
    """Remove all slides from a presentation, keeping layouts and masters."""
    r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    while len(prs.slides) > 0:
        sld_id = prs.slides._sldIdLst[0]
        r_id = sld_id.get(f"{r_ns}id")
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(sld_id)


def create_pptx_from_template(
    template_path: str | Path,
    output_path: str | Path,
    slides_content: list[dict[str, Any]],
) -> None:
    """Create a new PPTX from a template, preserving all styles.

    Args:
        template_path: Path to the template .pptx file.
        output_path: Path for the output .pptx file.
        slides_content: List of dicts, each with:
            - layout_index (int): Index of the slide layout to use.
            - layout_name (str, optional): Layout name (used if layout_index missing).
            - title (str, optional): Title text.
            - body (str, optional): Body/content text (single string).
            - bullets (list[str], optional): Body as bullet points.
            - notes (str, optional): Speaker notes.
            - code (str, optional): Code block to add below the title.
    """
    prs = Presentation(str(template_path))
    _delete_all_slides(prs)

    # Build a name→index lookup for layouts
    layout_by_name = {}
    for i, layout in enumerate(prs.slide_layouts):
        layout_by_name[layout.name.lower()] = i

    for slide_data in slides_content:
        # Resolve layout
        layout_idx = slide_data.get("layout_index")
        if layout_idx is None:
            name = slide_data.get("layout_name", "Title and Content").lower()
            layout_idx = layout_by_name.get(name, 1)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Title
        title_text = slide_data.get("title")
        if title_text and slide.shapes.title:
            slide.shapes.title.text = title_text

        # Body — either a single string or bullet list
        body_text = slide_data.get("body")
        bullets = slide_data.get("bullets")
        if body_text or bullets:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    if bullets:
                        set_body_bullets(ph, bullets)
                    else:
                        ph.text = body_text
                    break

        # Code block
        code_text = slide_data.get("code")
        if code_text:
            add_code_block(slide, code_text)

        # Speaker notes
        notes_text = slide_data.get("notes")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    prs.save(str(output_path))
    print(f"Created {output_path} with {len(slides_content)} slides")


def set_body_bullets(placeholder: Any, lines: list[str]) -> None:
    """Set placeholder text as bullet points, one per line.

    Inherits the placeholder's default paragraph/bullet formatting.
    """
    tf = placeholder.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.level = 0


def add_code_block(
    slide: Any,
    code_text: str,
    left: int = Inches(0.5),
    top: int = Inches(2),
    width: int = Inches(9),
    height: int = Inches(4.5),
    font_size: int = Pt(11),
    font_name: str = "Consolas",
    bg_color: RGBColor = RGBColor(0x1E, 0x1E, 0x1E),
    fg_color: RGBColor = RGBColor(0xD4, 0xD4, 0xD4),
) -> None:
    """Add a code block as a text box with monospace font and dark background."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    para = tf.paragraphs[0]
    para.text = code_text
    for run in para.runs:
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = fg_color


def copy_slide_background(source_slide: Any, target_slide: Any) -> None:
    """Copy background properties from a source slide to a target slide."""
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    source_bg = source_slide._element.find("p:bg", nsmap)
    if source_bg is None:
        return

    existing_bg = target_slide._element.find("p:bg", nsmap)
    if existing_bg is not None:
        target_slide._element.remove(existing_bg)

    new_bg = copy.deepcopy(source_bg)
    # cSld is the first child; background should be its first child
    csld = target_slide._element.find("p:cSld", nsmap)
    if csld is not None:
        csld.insert(0, new_bg)


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: uv run pptx_from_template.py <template.pptx>")
        print("  Inspects a PPTX template and prints its structure.")
        sys.exit(1)

    print_template_info(sys.argv[1])
